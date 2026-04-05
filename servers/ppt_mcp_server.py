"""
ppt_mcp_server.py — Stateful MCP server for PowerPoint generation.
Supports: Multiple themes, alternating slide layouts, image insertion.
All layout values are tuned for a 13.33 × 7.5 inch (widescreen) slide.
"""
import os
import sys
import io
import json
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import mcp.server.stdio
from mcp.server import Server
from mcp.types import Tool, TextContent

# Add parent dir to path so we can import templates
sys.path.insert(0, os.path.dirname(__file__))
from templates import get_template, TEMPLATES

def debug_print(msg):
    print(f"[DEBUG] {msg}", file=sys.stderr)
    sys.stderr.flush()

# ─── Global State ────────────────────────────────────────────────────────────
prs = None
current_theme = {}
slide_count = 0
IMAGE_CACHE_DIR = os.path.join(os.path.dirname(__file__), "..", "output", "img_cache")
os.makedirs(IMAGE_CACHE_DIR, exist_ok=True)

app = Server("ppt-mcp")

# ─── Layout Constants (inches, for 13.33 × 7.5 slide) ───────────────────────
MARGIN_LEFT   = 1.2      # generous left padding
MARGIN_RIGHT  = 1.2      # mirror the left
MARGIN_TOP    = 0.25      # breathing room above title bar
TITLE_BAR_H   = 1.4       # height of colored title banner
CONTENT_TOP   = 1.9       # where bullet text starts (below title bar + gap)
CONTENT_BOTTOM_PAD = 0.6  # padding at bottom
BULLET_FONT   = 20        # readable bullet size
TITLE_FONT    = 30        # slide title size
COVER_TITLE   = 44        # cover slide title
COVER_SUB     = 22        # cover slide subtitle
BULLET_SPACING_PT = 12    # space after each bullet paragraph
LINE_SPACING_PT   = 18    # line spacing within a paragraph

# ─── Helpers ─────────────────────────────────────────────────────────────────

def set_slide_background(slide, rgb_tuple):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_tuple)

def _add_rich_textbox(
    slide, text_lines, left, top, width, height,
    font_size=20, bold=False, color_rgb=(0,0,0),
    font_name="Calibri", align=PP_ALIGN.LEFT,
    space_after=Pt(12), line_spacing=None,
    bullet_char=None, anchor=MSO_ANCHOR.TOP
):
    """
    Creates a textbox with one paragraph per item in *text_lines*.
    Each paragraph gets its own run, spacing, and optional bullet character.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None          # let the box size stay fixed
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    for idx, line in enumerate(text_lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = align
        if space_after:
            p.space_after = space_after
        if line_spacing:
            p.line_spacing = line_spacing

        display = f"{bullet_char}  {line}" if bullet_char else line
        run = p.add_run()
        run.text = display
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color_rgb)
        run.font.name = font_name

    return txBox

def download_unsplash_image(query: str, index: int = 0):
    """Downloads a royalty-free image from Unsplash Source API."""
    try:
        url = f"https://source.unsplash.com/960x540/?{query.replace(' ', ',')}&sig={index}"
        debug_print(f"Fetching image for '{query}' …")
        resp = requests.get(url, timeout=10, stream=True)
        if resp.status_code == 200:
            img_path = os.path.join(IMAGE_CACHE_DIR, f"slide_img_{index}.jpg")
            with open(img_path, 'wb') as f:
                for chunk in resp.iter_content(1024):
                    f.write(chunk)
            debug_print(f"Image saved → {img_path}")
            return img_path
    except Exception as e:
        debug_print(f"Image download failed: {e}")
    return None

# ─── Slide Builders ──────────────────────────────────────────────────────────

def make_title_slide(prs_obj, title_text: str, subtitle_text: str, theme: dict):
    """Cover slide — large centred title, subtitle, accent bars."""
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[6])
    sw = prs_obj.slide_width.inches
    sh = prs_obj.slide_height.inches

    set_slide_background(slide, theme["bg_color"])

    # ── Top accent band ──────────────────────────────────────────────────
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(sw), Inches(3.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*theme["title_bg"])
    bar.line.fill.background()

    # ── Title ────────────────────────────────────────────────────────────
    _add_rich_textbox(
        slide, [title_text],
        left=1.5, top=0.6, width=sw - 3.0, height=2.2,
        font_size=COVER_TITLE, bold=True,
        color_rgb=theme["title_color"],
        font_name=theme["font_title"],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        space_after=None
    )

    # ── Subtitle ─────────────────────────────────────────────────────────
    _add_rich_textbox(
        slide, [subtitle_text],
        left=2.0, top=3.5, width=sw - 4.0, height=1.2,
        font_size=COVER_SUB, bold=False,
        color_rgb=theme["content_color"],
        font_name=theme["font_body"],
        align=PP_ALIGN.CENTER,
        space_after=None
    )

    # ── Bottom accent stripe ─────────────────────────────────────────────
    bot = slide.shapes.add_shape(1, Inches(0), Inches(sh - 0.35), Inches(sw), Inches(0.35))
    bot.fill.solid()
    bot.fill.fore_color.rgb = RGBColor(*theme["accent_color"])
    bot.line.fill.background()


def make_content_slide(prs_obj, title_text: str, bullets: list, theme: dict):
    """Full-width content slide with proper margins and bullet spacing."""
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[6])
    sw = prs_obj.slide_width.inches
    sh = prs_obj.slide_height.inches

    set_slide_background(slide, theme["bg_color"])

    # ── Title banner ─────────────────────────────────────────────────────
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(sw), Inches(TITLE_BAR_H)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(*theme["title_bg"])
    title_bar.line.fill.background()

    # ── Left accent stripe below banner ──────────────────────────────────
    stripe = slide.shapes.add_shape(
        1, Inches(0), Inches(TITLE_BAR_H),
        Inches(0.15), Inches(sh - TITLE_BAR_H)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(*theme["accent_color"])
    stripe.line.fill.background()

    # ── Title text (inside banner, left-padded) ──────────────────────────
    _add_rich_textbox(
        slide, [title_text],
        left=MARGIN_LEFT, top=MARGIN_TOP,
        width=sw - MARGIN_LEFT - MARGIN_RIGHT, height=TITLE_BAR_H - 0.3,
        font_size=TITLE_FONT, bold=True,
        color_rgb=theme["title_color"],
        font_name=theme["font_title"],
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
        space_after=None
    )

    # ── Bullet content ───────────────────────────────────────────────────
    content_h = sh - CONTENT_TOP - CONTENT_BOTTOM_PAD
    content_w = sw - MARGIN_LEFT - MARGIN_RIGHT

    _add_rich_textbox(
        slide, bullets,
        left=MARGIN_LEFT, top=CONTENT_TOP,
        width=content_w, height=content_h,
        font_size=BULLET_FONT, bold=False,
        color_rgb=theme["content_color"],
        font_name=theme["font_body"],
        align=PP_ALIGN.LEFT,
        space_after=Pt(BULLET_SPACING_PT),
        bullet_char="▸"
    )


def make_image_slide(prs_obj, title_text: str, bullets: list, img_path: str, theme: dict):
    """Two-column slide: bullets on left, image on right, balanced."""
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[6])
    sw = prs_obj.slide_width.inches
    sh = prs_obj.slide_height.inches

    set_slide_background(slide, theme["bg_color"])

    # ── Title banner (full width) ────────────────────────────────────────
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(sw), Inches(TITLE_BAR_H)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(*theme["title_bg"])
    title_bar.line.fill.background()

    _add_rich_textbox(
        slide, [title_text],
        left=MARGIN_LEFT, top=MARGIN_TOP,
        width=sw - MARGIN_LEFT - MARGIN_RIGHT, height=TITLE_BAR_H - 0.3,
        font_size=TITLE_FONT, bold=True,
        color_rgb=theme["title_color"],
        font_name=theme["font_title"],
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
        space_after=None
    )

    # ── Layout columns ───────────────────────────────────────────────────
    usable_w = sw - MARGIN_LEFT - MARGIN_RIGHT
    col_gap = 0.4
    text_col_w = usable_w * 0.55        # ~55 % for text
    img_col_w  = usable_w - text_col_w - col_gap   # ~45 % for image

    content_h = sh - CONTENT_TOP - CONTENT_BOTTOM_PAD

    # ── Left: bullets ────────────────────────────────────────────────────
    _add_rich_textbox(
        slide, bullets,
        left=MARGIN_LEFT, top=CONTENT_TOP,
        width=text_col_w, height=content_h,
        font_size=BULLET_FONT - 2, bold=False,  # slightly smaller to fit beside image
        color_rgb=theme["content_color"],
        font_name=theme["font_body"],
        align=PP_ALIGN.LEFT,
        space_after=Pt(BULLET_SPACING_PT),
        bullet_char="▸"
    )

    # ── Right: image ─────────────────────────────────────────────────────
    if img_path and os.path.exists(img_path):
        try:
            img_left  = Inches(MARGIN_LEFT + text_col_w + col_gap)
            img_top   = Inches(CONTENT_TOP)
            img_w     = Inches(img_col_w)
            img_h     = Inches(content_h)
            slide.shapes.add_picture(img_path, img_left, img_top, img_w, img_h)
            debug_print(f"Image inserted ✓")
        except Exception as e:
            debug_print(f"Image insert failed: {e}")

    # ── Bottom accent stripe ─────────────────────────────────────────────
    bot = slide.shapes.add_shape(
        1, Inches(0), Inches(sh - 0.2), Inches(sw), Inches(0.2)
    )
    bot.fill.solid()
    bot.fill.fore_color.rgb = RGBColor(*theme["accent_color"])
    bot.line.fill.background()


# ─── MCP Tool Registration ────────────────────────────────────────────────────

@app.list_tools()
async def list_tools() -> list[Tool]:
    return [
        Tool(
            name="initialize_presentation",
            description=(
                "Initializes a new PowerPoint presentation with the given theme and topic. "
                "Call this FIRST before any write_slide calls. "
                "Supported themes: modern_minimal, dark_tech, startup_pitch, "
                "education_visual, corporate_professional, academic_clean."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "topic": {"type": "string", "description": "The presentation topic (used in the title slide)."},
                    "template": {"type": "string", "description": "Theme name. One of: modern_minimal, dark_tech, startup_pitch, education_visual, corporate_professional, academic_clean."},
                    "slide_count": {"type": "integer", "description": "Target number of content slides to generate (3-10)."}
                },
                "required": ["topic", "template", "slide_count"]
            }
        ),
        Tool(
            name="write_slide",
            description=(
                "Adds a content slide with title, bullet points, and optionally an image. "
                "Automatically alternates between plain content slides and image slides. "
                "For image slides, provide a search query to fetch a relevant image."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Title of the slide."},
                    "bullets": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of 3-5 bullet points for the slide."
                    },
                    "image_query": {
                        "type": "string",
                        "description": "Optional: A search term to fetch a relevant image (e.g., 'artificial intelligence robot'). Leave empty for no image."
                    }
                },
                "required": ["title", "bullets"]
            }
        ),
        Tool(
            name="save_ppt",
            description="Saves the current in-memory presentation to disk. MUST be called at the end after all slides are added.",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Absolute path to save the .pptx file."}
                },
                "required": ["filename"]
            }
        )
    ]


@app.call_tool()
async def call_tool(name: str, arguments: dict) -> list[TextContent]:
    global prs, current_theme, slide_count

    if name == "initialize_presentation":
        topic = arguments.get("topic", "Presentation")
        template_name = arguments.get("template", "corporate_professional")
        slide_count = arguments.get("slide_count", 5)
        current_theme = get_template(template_name)

        debug_print(f"Init — Topic: {topic}, Theme: {template_name}, Slides: {slide_count}")

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        make_title_slide(prs, topic, f"A {slide_count}-slide presentation", current_theme)

        return [TextContent(type="text", text=f"Presentation initialized with theme '{template_name}'. Title slide created. Now add {slide_count} content slides using write_slide.")]

    elif name == "write_slide":
        if prs is None:
            return [TextContent(type="text", text="Error: Call initialize_presentation first.")]

        title_text = arguments.get("title", "Slide")
        bullets = arguments.get("bullets", [])
        image_query = arguments.get("image_query", "")

        debug_print(f"Slide: {title_text} | img_query: '{image_query}'")

        idx = len(prs.slides)
        if image_query and idx % 2 == 0:
            img_path = download_unsplash_image(image_query, index=idx)
            make_image_slide(prs, title_text, bullets, img_path, current_theme)
            return [TextContent(type="text", text=f"Added image slide '{title_text}'.")]
        else:
            make_content_slide(prs, title_text, bullets, current_theme)
            return [TextContent(type="text", text=f"Added content slide '{title_text}'.")]

    elif name == "save_ppt":
        if prs is None:
            return [TextContent(type="text", text="Error: No presentation to save.")]

        filename = arguments.get("filename", r"C:\Users\somas\Downloads\ppt_antigravity\output.pptx")
        if not os.path.isabs(filename):
            filename = os.path.join(os.path.dirname(os.path.dirname(__file__)), filename)

        debug_print(f"Saving → {filename}")
        prs.save(filename)

        exists = os.path.exists(filename)
        size = os.path.getsize(filename) if exists else 0
        debug_print(f"Exists: {exists}, Size: {size} bytes")

        if not exists:
            return [TextContent(type="text", text=f"ERROR: File NOT created at {filename}")]

        return [TextContent(type="text", text=f"SUCCESS: Saved to {filename} ({size} bytes). File verified.")]
    else:
        raise ValueError(f"Unknown tool: {name}")


async def main():
    async with mcp.server.stdio.stdio_server() as (read_stream, write_stream):
        await app.run(read_stream, write_stream, app.create_initialization_options())

if __name__ == "__main__":
    import asyncio
    asyncio.run(main())
